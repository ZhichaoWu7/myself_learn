import os
import re
from pathlib import Path
import pymupdf4llm
from docx import Document
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE

class DocToMdConverter:
    def __init__(self, out_dir='converted_md', image_dir='extracted_image'):
        # 保持你的绝对路径初始化
        self.out_dir = Path(out_dir).absolute()
        self.image_dir = Path(image_dir).absolute()
        self.out_dir.mkdir(exist_ok=True, parents=True)
        self.image_dir.mkdir(exist_ok=True, parents=True)

    def convert_and_replace(self, file_path: str):
        p = Path(file_path)
        if not p.exists():
            print(f"❌ 文件不存在: {file_path}")
            return None

        content = self.convert(p)

        if content:
            clean_content = content.replace('\x00', '')
            md_file_path = self.out_dir / f"{p.stem}.md"

            try:
                # 使用 utf-8-sig 写入，这是解决很多编辑器打开乱码的“特效药”
                with open(md_file_path, 'w', encoding='utf-8-sig', errors='replace') as f:
                    f.write(clean_content)

                print(f"✅ 转换成功并保存: {md_file_path.name}")
                print(f"🗑️ 正在删除原始文件: {p.name}")
                os.remove(p)
                return str(md_file_path)
            except Exception as e:
                print(f"❌ 写入文件或删除原文件时出错: {e}")
        else:
            print(f"⚠️ 文件 {p.name} 转换后内容为空，跳过删除。")
        return None

    def convert(self, file_path: Path) -> str:
        ext = file_path.suffix.lower()
        print(f"🚀 正在顺序解析: {file_path.name}")

        if ext in [".docx", ".doc"]:
            return self._convert_doc(file_path)
        elif ext in [".ppt", ".pptx"]:
            return self._convert_ppt(file_path)
        elif ext == ".pdf":
            return self._convert_pdf(file_path)
        elif ext in [".txt", ".md"]:
            return self._read_text_safe(file_path)
        return ""

    def _convert_pdf(self, path: Path) -> str:
        """极简稳健版：不折腾，能跑通就行"""
        try:
            # 1. 尝试直接转换
            md_text = pymupdf4llm.to_markdown(
                str(path),
                write_images=True,
                image_path=str(self.image_dir),
                image_format="png",
                dpi=150  # 稍微调低，防止‘图带字’占空间
            )

            # 2. 路径统一重定向到 picture/
            abs_prefix = str(self.image_dir)
            md_text = md_text.replace(abs_prefix, "picture")

            # 3. 简单的乱码‘物理清洗’
            # 过滤掉那些不可见的控制字符，减少 LLM 的困扰
            md_text = "".join(ch for ch in md_text if ch.isprintable() or ch in "\n\t")

            return md_text
        except:
            return "PDF 解析失败，建议手动检查。"

    def _read_text_safe(self, path: Path) -> str:
        for enc in ['utf-8', 'gbk', 'utf-16']:
            try:
                with open(path, 'r', encoding=enc) as f:
                    return f.read()
            except UnicodeDecodeError:
                continue
        return ""

    def _convert_doc(self, path: Path) -> str:
        doc = Document(str(path))
        md_lines = []
        img_counter = 0
        images_map = {rel.rId: rel.target_part.blob for rel in doc.part.rels.values() if "image" in rel.target_ref}

        for para in doc.paragraphs:
            if para.text.strip():
                md_lines.append(para.text.strip())
            rIds = re.findall(r'r:embed="(rId\d+)"', para._p.xml)
            for rId in rIds:
                if rId in images_map:
                    img_counter += 1
                    img_name = f"{path.stem}_docx_img_{img_counter}.png"
                    img_save_path = self.image_dir / img_name
                    with open(img_save_path, "wb") as f:
                        f.write(images_map[rId])
                    # 💡 修正：使用相对路径 picture/ 保证预览可见
                    md_lines.append(f"\n![img_{img_counter}](picture/{img_name})\n")
        return "\n".join(md_lines)

    def _convert_ppt(self, path: Path) -> str:
        prs = Presentation(str(path))
        md_lines = [f"# PPT: {path.name}\n"]
        img_counter = 0
        for i, slide in enumerate(prs.slides):
            md_lines.append(f"## Slide {i + 1}")
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text.strip():
                    md_lines.append(shape.text.strip())
                if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                    img_counter += 1
                    image = shape.image
                    img_name = f"{path.stem}_pptx_s{i + 1}_img_{img_counter}.{image.ext}"
                    img_save_path = self.image_dir / img_name
                    with open(img_save_path, "wb") as f:
                        f.write(image.blob)
                    # 💡 修正：使用相对路径 picture/ 保证预览可见
                    md_lines.append(f"\n![img_{img_counter}](picture/{img_name})\n")
            md_lines.append("\n---\n")
        return "\n".join(md_lines)